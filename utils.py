import fitz
from pptx import Presentation

def extract_pdf_text(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def export_slides(slides, output_file="slides.pptx"):
    prs = Presentation()
    for content in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        lines = content.split('\n')
        title = lines[0].replace("Slide Title:", "").strip()
        bullets = [line.strip("-• ") for line in lines[1:] if line.strip()]
        slide.shapes.title.text = title
        content_frame = slide.shapes.placeholders[1].text_frame
        content_frame.clear()
        for bullet in bullets:
            content_frame.add_paragraph().text = bullet
    prs.save(output_file)
    return output_file
